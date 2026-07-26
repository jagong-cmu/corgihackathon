"""Turning an uploaded file into text.

The first stage of the ingestion path: bytes in, plain text out, ready for
`documents.chunk_document`. Everything downstream — chunking, embedding,
`show_source` on the board — only ever sees the string this module returns.

Format support is deliberately narrow. A learner uploads slides, a problem set,
or lecture notes; that is pdf/docx/pptx/text and nothing else. An unsupported
file gets a named error rather than a best-effort decode, because a PowerPoint
run through a text decoder produces chunks made of XML fragments that retrieve
plausibly and read as garbage the moment `show_source` puts one on the board in
front of the learner.

The parsers are an optional extra so the offline test suite keeps installing
with no dependencies:

    uv sync --extra documents
"""

from __future__ import annotations

import io
import logging
from dataclasses import dataclass
from pathlib import PurePosixPath

log = logging.getLogger(__name__)

# Extension -> kind. Content types are checked first when present, but browsers
# lie about them often enough (docx as application/zip, md as text/plain) that
# the extension is the tiebreak rather than the other way round.
_BY_EXTENSION = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".pptx": "pptx",
    ".txt": "text",
    ".md": "text",
    ".markdown": "text",
    ".csv": "text",
    ".rst": "text",
    ".json": "text",
}

_BY_CONTENT_TYPE = {
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
}

SUPPORTED_EXTENSIONS = tuple(sorted(_BY_EXTENSION))

# The legacy binary formats. Called out separately because "we can't read .doc"
# is a much more useful message than "unsupported file", and the fix (re-save as
# .docx) is something the learner can actually do.
_LEGACY = {".doc": "Word 97-2003", ".ppt": "PowerPoint 97-2003", ".xls": "Excel 97-2003"}


class UnsupportedDocumentError(ValueError):
    """The file is not a format we can turn into teachable text."""


class ExtractionDependencyError(RuntimeError):
    """The format is supported but its parser is not installed."""


@dataclass(frozen=True)
class ExtractedDocument:
    text: str
    kind: str
    """One of pdf | docx | pptx | text."""

    pages: int | None = None
    """Pages, slides, or None for flat text. Recorded on the upload so the UI can
    say '12 slides' rather than a byte count."""


def detect_kind(filename: str, content_type: str | None = None) -> str:
    """Classify an upload, or raise `UnsupportedDocumentError` naming why."""
    suffix = PurePosixPath(filename or "").suffix.lower()

    if content_type:
        # Strip any '; charset=' parameter before matching.
        base = content_type.split(";", 1)[0].strip().lower()
        if base in _BY_CONTENT_TYPE:
            return _BY_CONTENT_TYPE[base]
        if base.startswith("text/") and suffix not in _BY_EXTENSION:
            # text/* with an extension we don't know: trust the content type.
            return "text"

    if suffix in _BY_EXTENSION:
        return _BY_EXTENSION[suffix]

    if suffix in _LEGACY:
        raise UnsupportedDocumentError(
            f"{filename} is a {_LEGACY[suffix]} file. Re-save it as "
            f"{suffix}x and upload that."
        )

    raise UnsupportedDocumentError(
        f"cannot extract text from {filename or 'the upload'}"
        f" (content type {content_type or 'unknown'}). Supported: "
        + ", ".join(SUPPORTED_EXTENSIONS)
    )


def extract_text(
    data: bytes, *, filename: str, content_type: str | None = None
) -> ExtractedDocument:
    """Extract text from one uploaded file.

    Raises `UnsupportedDocumentError` for a format we don't handle and
    `ExtractionDependencyError` when the parser for a supported format is
    missing. Both are the caller's problem to report — an empty string would
    ingest cleanly and produce a document the tutor silently knows nothing about.
    """
    kind = detect_kind(filename, content_type)
    if kind == "pdf":
        return _extract_pdf(data)
    if kind == "docx":
        return _extract_docx(data)
    if kind == "pptx":
        return _extract_pptx(data)
    return _extract_text(data)


def _extract_text(data: bytes) -> ExtractedDocument:
    # errors="replace" rather than strict: a stray byte in an otherwise fine
    # notes file should cost one character, not the whole upload.
    return ExtractedDocument(text=data.decode("utf-8", errors="replace"), kind="text")


def _extract_pdf(data: bytes) -> ExtractedDocument:
    try:
        from pypdf import PdfReader
    except ImportError as exc:  # pragma: no cover - exercised by the error path test
        raise ExtractionDependencyError(
            "reading PDFs needs pypdf — install the extra: uv sync --extra documents"
        ) from exc

    reader = PdfReader(io.BytesIO(data))
    # Pages are joined with a blank line because chunk_document splits on
    # paragraph boundaries first, and a page break is a real boundary.
    pages = [page.extract_text() or "" for page in reader.pages]
    text = "\n\n".join(p.strip() for p in pages if p.strip())

    if not text.strip():
        # A scanned PDF is images with no text layer. Say so: the alternative is
        # an upload that reports success and then answers nothing.
        raise UnsupportedDocumentError(
            "this PDF has no extractable text — it looks like a scan. "
            "OCR it first, or upload the original document."
        )
    return ExtractedDocument(text=text, kind="pdf", pages=len(reader.pages))


def _extract_docx(data: bytes) -> ExtractedDocument:
    try:
        import docx
    except ImportError as exc:  # pragma: no cover
        raise ExtractionDependencyError(
            "reading .docx needs python-docx — install the extra: uv sync --extra documents"
        ) from exc

    document = docx.Document(io.BytesIO(data))
    blocks = [p.text.strip() for p in document.paragraphs]

    # Tables carry the content in a lot of problem sets, and python-docx does not
    # include them in `paragraphs`. Rows become one line each so a row still
    # reads as a unit after chunking.
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                blocks.append(" | ".join(cells))

    return ExtractedDocument(text="\n\n".join(b for b in blocks if b), kind="docx")


def _extract_pptx(data: bytes) -> ExtractedDocument:
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover
        raise ExtractionDependencyError(
            "reading .pptx needs python-pptx — install the extra: uv sync --extra documents"
        ) from exc

    presentation = Presentation(io.BytesIO(data))
    slides: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts = [
            shape.text_frame.text.strip()
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        if not parts:
            continue
        # The slide number is kept in the text on purpose: a learner asking
        # "what was on slide 6" is asking a question the index can only answer
        # if the number survived extraction.
        slides.append(f"Slide {index}\n" + "\n".join(parts))

    return ExtractedDocument(
        text="\n\n".join(slides), kind="pptx", pages=len(presentation.slides)
    )
